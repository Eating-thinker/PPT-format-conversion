import streamlit as st
from pptx import Presentation
import subprocess
import tempfile

# -----------------------------
# 預設風格選單
# -----------------------------
PRESET_STYLES = {
    "科技感黑銀": "科技感、黑銀配色、俐落線條、霧面金屬風格、未來介面 UI。",
    "簡約蘋果風": "極簡留白、蘋果風格、柔和陰影、高質感黑白灰。",
    "商務專業藍": "企業藍、正式商務、整齊結構化排版、乾淨專業。",
    "溫暖奶油風": "米白色、低彩度、柔和圓角、溫暖質感、療癒風系。",
    "活潑卡通風": "明亮色彩、卡通插畫風、大圖示與活潑字體。"
}

# 初始化 session_state（避免風格選擇消失）
if "selected_style" not in st.session_state:
    st.session_state.selected_style = ""

# -----------------------------
# function：讀 PPT
# -----------------------------
def extract_text_from_ppt(path):
    prs = Presentation(path)
    result = []

    for slide in prs.slides:
        items = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                items.append(shape.text)
        result.append("\n".join(items))

    return result


# -----------------------------
# function：用本地模型（Ollama）產生設計
# -----------------------------
def generate_redesign(slides_text, style):

    prompt = f"""
你是一位頂尖簡報設計師。

以下是原始 PPT 的內容：
{slides_text}

請依照以下風格重新設計：
{style}

請輸出下面格式：

[Slide 1 Title]
文字
[Slide 1 Bullets]
- 一
- 二
- 三

[Slide 2 Title]
...
（所有頁面依序輸出）
    """

    process = subprocess.Popen(
        ["ollama", "run", "llama3.1"],
        stdin=subprocess.PIPE,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True
    )

    output, _ = process.communicate(prompt)

    return output


# -----------------------------
# function：生成新 PPT
# -----------------------------
def create_new_ppt(design_text):
    prs = Presentation()

    slides = design_text.split("[Slide")
    slides = [s.strip() for s in slides if s.strip()]

    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        lines = s.split("\n")

        if len(lines) < 3:
            continue

        # 標題
        title = slide.shapes.title
        title.text = lines[1].strip()

        # 內容
        body = slide.placeholders[1].text_frame
        for l in lines[2:]:
            if l.strip().startswith("-"):
                p = body.add_paragraph()
                p.text = l.replace("-", "").strip()
            else:
                p = body.add_paragraph()
                p.text = l.strip()

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name


# -----------------------------
# Streamlit UI
# -----------------------------
st.set_page_config(page_title="免費版 AI PPT 設計", page_icon="🎨")

st.title("🎨 免費版 AI PPT 重新設計（本地模型，不用 API Key）")
st.subheader("使用 Ollama + 本地模型（LLaMA3 / Qwen）完全免費")

# -----------------------------
# 風格選擇
# -----------------------------
st.markdown("### 🔥 選擇風格")

cols = st.columns(5)
keys = list(PRESET_STYLES.keys())

for i, col in enumerate(cols):
    with col:
        # 修正：按下按钮後，將值儲存到 session_state
        if st.button(keys[i]):
            st.session_state.selected_style = PRESET_STYLES[keys[i]]

# 顯示選到的風格
if st.session_state.selected_style:
    st.success(f"已選擇風格：{st.session_state.selected_style}")

# 自訂風格
custom_style = st.text_area("🖋 自訂風格（可留空）")

# 上傳 PPT
uploaded = st.file_uploader("📤 上傳 PPT (.pptx)", type=["pptx"])

# -----------------------------
# 開始轉換
# -----------------------------
if st.button("🚀 開始轉換"):

    if not uploaded:
        st.error("請先上傳 PPT")
        st.stop()

    style = custom_style if custom_style else st.session_state.selected_style

    if not style:
        st.error("請先選擇或輸入風格")
        st.stop()

    # 保存原始 PPT 到暫存
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    tmp.write(uploaded.read())

    with st.spinner("讀取原始 PPT..."):
        content = extract_text_from_ppt(tmp.name)

    with st.spinner("AI 設計中（使用本地模型）..."):
        design = generate_redesign(content, style)

    with st.spinner("正在生成新 PPT…"):
        output_path = create_new_ppt(design)

    st.success("✨ 完成！PPT 已重新設計")
    st.download_button(
        "📥 下載全新的 PPT",
        data=open(output_path, "rb").read(),
        file_name="AI_redesigned_free_version.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
